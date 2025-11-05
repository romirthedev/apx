"""
Document Manager Plugin - Create and manage various document types
"""

import os
import subprocess
import logging
from datetime import datetime
from typing import Dict, Any, Optional

logger = logging.getLogger(__name__)

class DocumentManager:
    def __init__(self, google_sheets_manager=None):
        self.supported_formats = [
            'docx', 'xlsx', 'pptx', 'pdf', 'txt', 'md', 'csv', 'json'
        ]
        self.google_sheets_manager = google_sheets_manager
    
    def create_word_document(self, filename: str, content: str = "") -> str:
        """Create a new Word document."""
        try:
            # For macOS, we can use AppleScript to create documents
            if not filename.endswith('.docx'):
                filename += '.docx'
            
            full_path = os.path.expanduser(f"~/Documents/{filename}")
            
            # Create basic DOCX using python-docx if available
            try:
                from docx import Document
                doc = Document()
                if content:
                    doc.add_paragraph(content)
                else:
                    doc.add_paragraph(f"Document created on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
                doc.save(full_path)
                
                # Open the document
                subprocess.run(['open', full_path], check=True)
                return f"✅ Created and opened Word document: {filename}"
                
            except ImportError:
                # Fallback: create RTF file that can be opened by Word
                rtf_content = f"""{{\\rtf1\\ansi\\deff0 {{\\fonttbl {{\\f0 Times New Roman;}}}}
\\f0\\fs24 {content or f"Document created on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"}
}}"""
                rtf_path = full_path.replace('.docx', '.rtf')
                with open(rtf_path, 'w') as f:
                    f.write(rtf_content)
                
                subprocess.run(['open', rtf_path], check=True)
                return f"✅ Created and opened document: {os.path.basename(rtf_path)}"
                
        except Exception as e:
            logger.error(f"Failed to create Word document: {str(e)}")
            return f"❌ Failed to create Word document: {str(e)}"
    
    def create_excel_spreadsheet(self, filename: str, data: Optional[Dict] = None) -> str:
        """Create a new Excel spreadsheet."""
        try:
            if not filename.endswith('.xlsx'):
                filename += '.xlsx'
            
            full_path = os.path.expanduser(f"~/Documents/{filename}")
            
            try:
                from openpyxl import Workbook
                wb = Workbook()
                ws = wb.active
                ws.title = "Sheet1"
                
                if data:
                    # Add data if provided
                    if isinstance(data, dict) and 'headers' in data:
                        # Add headers
                        for col, header in enumerate(data['headers'], 1):
                            ws.cell(row=1, column=col, value=header)
                        
                        # Add rows
                        if 'rows' in data:
                            for row_idx, row_data in enumerate(data['rows'], 2):
                                for col_idx, cell_value in enumerate(row_data, 1):
                                    ws.cell(row=row_idx, column=col_idx, value=cell_value)
                else:
                    # Add default content
                    ws['A1'] = 'Created Date'
                    ws['B1'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                
                wb.save(full_path)
                subprocess.run(['open', full_path], check=True)
                return f"✅ Created and opened Excel spreadsheet: {filename}"
                
            except ImportError:
                # Fallback: create CSV file
                csv_path = full_path.replace('.xlsx', '.csv')
                with open(csv_path, 'w') as f:
                    f.write(f"Created Date,{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\\n")
                    if data and 'headers' in data:
                        f.write(','.join(data['headers']) + '\\n')
                        if 'rows' in data:
                            for row in data['rows']:
                                f.write(','.join(str(cell) for cell in row) + '\\n')
                
                subprocess.run(['open', csv_path], check=True)
                return f"✅ Created and opened CSV file: {os.path.basename(csv_path)}"
                
        except Exception as e:
            logger.error(f"Failed to create Excel spreadsheet: {str(e)}")
            return f"❌ Failed to create Excel spreadsheet: {str(e)}"
    
    def create_presentation(self, filename: str, title: str = "", slides: Optional[list] = None) -> str:
        """Create a new PowerPoint presentation."""
        try:
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            full_path = os.path.expanduser(f"~/Documents/{filename}")
            
            try:
                from pptx import Presentation
                prs = Presentation()
                
                # Title slide
                title_slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(title_slide_layout)
                slide.shapes.title.text = title or "New Presentation"
                slide.placeholders[1].text = f"Created on {datetime.now().strftime('%Y-%m-%d')}"
                
                # Add additional slides if provided
                if slides:
                    bullet_slide_layout = prs.slide_layouts[1]
                    for slide_data in slides:
                        slide = prs.slides.add_slide(bullet_slide_layout)
                        slide.shapes.title.text = slide_data.get('title', 'Slide Title')
                        if 'content' in slide_data:
                            slide.placeholders[1].text = slide_data['content']
                
                prs.save(full_path)
                subprocess.run(['open', full_path], check=True)
                return f"✅ Created and opened PowerPoint presentation: {filename}"
                
            except ImportError:
                # Fallback: create a text outline
                txt_path = full_path.replace('.pptx', '_outline.txt')
                with open(txt_path, 'w') as f:
                    f.write(f"PRESENTATION OUTLINE: {title or 'New Presentation'}\\n")
                    f.write(f"Created: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\\n\\n")
                    
                    if slides:
                        for i, slide_data in enumerate(slides, 1):
                            f.write(f"SLIDE {i}: {slide_data.get('title', 'Slide Title')}\\n")
                            if 'content' in slide_data:
                                f.write(f"  - {slide_data['content']}\\n")
                            f.write("\\n")
                
                subprocess.run(['open', txt_path], check=True)
                return f"✅ Created presentation outline: {os.path.basename(txt_path)}"
                
        except Exception as e:
            logger.error(f"Failed to create presentation: {str(e)}")
            return f"❌ Failed to create presentation: {str(e)}"
    
    def create_pdf_document(self, filename: str, content: str = "") -> str:
        """Create a new PDF document."""
        try:
            if not filename.endswith('.pdf'):
                filename += '.pdf'
            
            full_path = os.path.expanduser(f"~/Documents/{filename}")
            
            try:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                
                c = canvas.Canvas(full_path, pagesize=letter)
                width, height = letter
                
                # Add title
                c.setFont("Helvetica-Bold", 16)
                c.drawString(72, height - 72, filename.replace('.pdf', ''))
                
                # Add content
                c.setFont("Helvetica", 12)
                text_content = content or f"PDF document created on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
                
                # Simple text wrapping
                lines = text_content.split('\\n')
                y_position = height - 120
                for line in lines:
                    if y_position < 72:  # Start new page if needed
                        c.showPage()
                        y_position = height - 72
                    c.drawString(72, y_position, line)
                    y_position -= 20
                
                c.save()
                subprocess.run(['open', full_path], check=True)
                return f"✅ Created and opened PDF document: {filename}"
                
            except ImportError:
                # Fallback: create HTML that can be printed to PDF
                html_path = full_path.replace('.pdf', '.html')
                html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <title>{filename.replace('.pdf', '')}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 40px; }}
        h1 {{ color: #333; }}
        .content {{ line-height: 1.6; }}
    </style>
</head>
<body>
    <h1>{filename.replace('.pdf', '')}</h1>
    <div class="content">
        <p><strong>Created:</strong> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
        <p>{content or 'This document was created by Apex.'}</p>
    </div>
</body>
</html>
"""
                with open(html_path, 'w') as f:
                    f.write(html_content)
                
                subprocess.run(['open', html_path], check=True)
                return f"✅ Created HTML document: {os.path.basename(html_path)} (You can print this to PDF)"
                
        except Exception as e:
            logger.error(f"Failed to create PDF document: {str(e)}")
            return f"❌ Failed to create PDF document: {str(e)}"
    
    def get_document_templates(self) -> Dict[str, Any]:
        """Get available document templates."""
        return {
            'word': {
                'letter': 'Business letter template',
                'resume': 'Professional resume template',
                'report': 'Business report template',
                'memo': 'Office memorandum template'
            },
            'excel': {
                'budget': 'Personal/business budget template',
                'invoice': 'Invoice template',
                'timesheet': 'Time tracking template',
                'inventory': 'Inventory management template'
            },
            'powerpoint': {
                'business': 'Business presentation template',
                'pitch': 'Startup pitch deck template',
                'training': 'Training presentation template',
                'report': 'Status report template'
            }
        }
    
    def create_from_template(self, doc_type: str, template: str, filename: str, **kwargs) -> str:
        """Create a document from a template."""
        templates = self.get_document_templates()
        
        if doc_type not in templates or template not in templates[doc_type]:
            return f"❌ Template '{template}' not found for document type '{doc_type}'"
        
        try:
            if doc_type == 'word':
                return self._create_word_template(template, filename, **kwargs)
            elif doc_type == 'excel':
                return self._create_excel_template(template, filename, **kwargs)
            elif doc_type == 'powerpoint':
                return self._create_ppt_template(template, filename, **kwargs)
            else:
                return f"❌ Unsupported document type: {doc_type}"
        except Exception as e:
            return f"❌ Failed to create document from template: {str(e)}"
    
    def _create_word_template(self, template: str, filename: str, **kwargs) -> str:
        """Create Word document from template."""
        content_map = {
            'letter': f"""[Your Name]
[Your Address]
[City, State ZIP Code]
[Email Address]
[Phone Number]

{datetime.now().strftime('%B %d, %Y')}

[Recipient Name]
[Title]
[Company Name]
[Address]

Dear [Recipient Name],

[Letter content goes here...]

Sincerely,
[Your Name]""",
            'resume': f"""[Your Full Name]
[Your Address] | [Phone] | [Email] | [LinkedIn]

PROFESSIONAL SUMMARY
[Brief summary of your experience and skills]

EXPERIENCE
[Job Title] - [Company Name] ([Start Date] - [End Date])
• [Achievement or responsibility]
• [Achievement or responsibility]

EDUCATION
[Degree] - [University Name] ([Graduation Year])

SKILLS
• [Skill 1]
• [Skill 2]
• [Skill 3]""",
            'report': f"""BUSINESS REPORT

Title: [Report Title]
Date: {datetime.now().strftime('%B %d, %Y')}
Prepared by: [Your Name]

EXECUTIVE SUMMARY
[Brief overview of the report]

INTRODUCTION
[Background and purpose]

FINDINGS
[Key findings and analysis]

RECOMMENDATIONS
[Action items and next steps]

CONCLUSION
[Summary and final thoughts]""",
            'memo': f"""MEMORANDUM

TO: [Recipient]
FROM: [Your Name]
DATE: {datetime.now().strftime('%B %d, %Y')}
RE: [Subject]

[Memo content goes here...]

[Your Name]"""
        }
        
        content = content_map.get(template, "Template content not available")
        return self.create_word_document(filename, content)
    
    def _create_excel_template(self, template: str, filename: str, **kwargs) -> str:
        """Create Excel document from template."""
        template_data = {
            'budget': {
                'headers': ['Category', 'Budgeted', 'Actual', 'Difference'],
                'rows': [
                    ['Income', '5000', '4800', '-200'],
                    ['Housing', '1500', '1500', '0'],
                    ['Food', '600', '650', '50'],
                    ['Transportation', '400', '380', '-20'],
                    ['Utilities', '200', '220', '20'],
                    ['Entertainment', '300', '280', '-20']
                ]
            },
            'invoice': {
                'headers': ['Item', 'Quantity', 'Rate', 'Amount'],
                'rows': [
                    ['Service/Product 1', '1', '100.00', '100.00'],
                    ['Service/Product 2', '2', '75.00', '150.00'],
                    ['', '', 'Total:', '250.00']
                ]
            },
            'timesheet': {
                'headers': ['Date', 'Start Time', 'End Time', 'Break', 'Total Hours', 'Project'],
                'rows': [
                    [datetime.now().strftime('%Y-%m-%d'), '9:00 AM', '5:00 PM', '1 hour', '7', 'Project A']
                ]
            }
        }
        
        data = template_data.get(template, None)
        return self.create_excel_spreadsheet(filename, data)
    
    def create_stock_document(self, company: str, data_type: str = "stock", period: str = "daily", format: str = "xlsx") -> str:
        """Create a document with stock data for a company.
        
        Args:
            company: Company name or ticker symbol
            data_type: Type of financial data (stock, earnings, revenue, etc.)
            period: Time period (daily, weekly, monthly, annual, quarterly)
            format: Document format (xlsx, docx, pdf)
            
        Returns:
            Status message
        """
        try:
            if not self.google_sheets_manager:
                return "❌ Cannot create stock document: Google Sheets Manager not available"
            
            # Fetch financial data using the Google Sheets Manager
            financial_data = self.google_sheets_manager.fetch_financial_data(company, data_type, period)
            
            if not financial_data["headers"] or not financial_data["rows"]:
                return f"❌ Failed to fetch financial data for {company}"
            
            # Create a title for the document
            title = f"{company.title()} {data_type.title()} {period.title()} - {datetime.now().strftime('%Y-%m-%d')}"
            
            # Create document based on requested format
            if format.lower() in ['xlsx', 'excel', 'spreadsheet']:
                return self.create_excel_spreadsheet(title, financial_data)
            elif format.lower() in ['docx', 'word', 'document']:
                # Format data as text for Word document
                content = self._format_financial_data_as_text(title, financial_data)
                return self.create_word_document(title, content)
            elif format.lower() in ['pdf']:
                # Format data as text for PDF document
                content = self._format_financial_data_as_text(title, financial_data)
                return self.create_pdf_document(title, content)
            else:
                return f"❌ Unsupported format: {format}"
                
        except Exception as e:
            logger.error(f"Failed to create stock document: {str(e)}")
            return f"❌ Failed to create stock document: {str(e)}"
    
    def _format_financial_data_as_text(self, title: str, data: Dict) -> str:
        """Format financial data as text for documents."""
        content = f"{title}\n\n"
        
        # Add headers
        if "headers" in data and data["headers"]:
            content += " | ".join(data["headers"]) + "\n"
            content += "-" * (sum(len(h) for h in data["headers"]) + (3 * (len(data["headers"]) - 1))) + "\n"
        
        # Add rows
        if "rows" in data and data["rows"]:
            for row in data["rows"]:
                content += " | ".join(str(cell) for cell in row) + "\n"
        
        content += "\nDocument generated on " + datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        return content
    
    def _create_ppt_template(self, template: str, filename: str, **kwargs) -> str:
        """Create PowerPoint document from template."""
        template_slides = {
            'business': [
                {'title': 'Agenda', 'content': '• Introduction\\n• Current Status\\n• Key Points\\n• Next Steps'},
                {'title': 'Current Status', 'content': '• Key metrics\\n• Progress update\\n• Challenges'},
                {'title': 'Next Steps', 'content': '• Action items\\n• Timeline\\n• Responsibilities'}
            ],
            'pitch': [
                {'title': 'Problem', 'content': 'What problem are we solving?'},
                {'title': 'Solution', 'content': 'How do we solve it?'},
                {'title': 'Market', 'content': 'Who are our customers?'},
                {'title': 'Business Model', 'content': 'How do we make money?'},
                {'title': 'Ask', 'content': 'What do we need?'}
            ],
            'training': [
                {'title': 'Learning Objectives', 'content': '• Objective 1\\n• Objective 2\\n• Objective 3'},
                {'title': 'Module 1', 'content': 'Content for first module'},
                {'title': 'Module 2', 'content': 'Content for second module'},
                {'title': 'Summary', 'content': 'Key takeaways and next steps'}
            ]
        }
        
        title = kwargs.get('title', f"{template.title()} Presentation")
        slides = template_slides.get(template, [])
        return self.create_presentation(filename, title, slides)
